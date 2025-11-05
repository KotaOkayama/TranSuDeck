"""
PPTX生成モジュール
マークダウン形式のテキストをPowerPointプレゼンテーションに変換
"""
import logging
import re
from typing import List, Tuple, Union

from pptx import Presentation
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


class MarkdownParser:
    """マークダウンテキストを解析してPowerPoint用の構造化データに変換"""

    @staticmethod
    def parse_line(line: str) -> Tuple[str, int, dict]:
        """
        1行のマークダウンテキストを解析

        Returns:
            Tuple[str, int, dict]: (テキスト内容, インデントレベル, 書式情報)
        """
        if not line:
            return "", 0, {}

        indent_level = 0
        format_info = {
            "bold": False,
            "italic": False,
            "underline": False,
            "is_heading": False,
            "heading_level": 0,
            "is_bullet": False,
            "is_numbered": False,
            "number": None,
        }

        # インデント検出（スペースまたはタブ）
        indent_match = re.match(r"^(\s+)", line)
        if indent_match:
            indent_level = len(indent_match.group(1)) // 2  # 2スペース = 1レベル
            line = line.lstrip()

        # 見出し検出
        heading_match = re.match(r"^(#{1,6})\s+(.+)$", line)
        if heading_match:
            format_info["is_heading"] = True
            format_info["heading_level"] = len(heading_match.group(1))
            line = heading_match.group(2)
            return line, indent_level, format_info

        # 箇条書き検出（-のみ）
        bullet_match = re.match(r"^-\s+(.+)$", line)
        if bullet_match:
            format_info["is_bullet"] = True
            line = bullet_match.group(1)
            return line, indent_level, format_info

        # 番号付きリスト検出
        numbered_match = re.match(r"^(\d+)\.\s+(.+)$", line)
        if numbered_match:
            format_info["is_numbered"] = True
            format_info["number"] = int(numbered_match.group(1))
            line = numbered_match.group(2)
            return line, indent_level, format_info

        return line, indent_level, format_info

    @staticmethod
    def extract_formatting_segments(text: str) -> List[Tuple[str, bool, bool, bool]]:
        """
        テキストから書式情報を抽出してセグメントに分割

        Returns:
            List[Tuple[str, bool, bool, bool]]: (テキスト, 太字, 斜体, 下線)のリスト
        """
        if not text:
            return [("", False, False, False)]

        segments = []
        remaining_text = text

        # 書式パターンを順番に処理
        while remaining_text:
            # 太字+斜体: ***text*** or ___text___
            match = re.match(r"^(.*?)(\*\*\*(.+?)\*\*\*|___(.+?)___)", remaining_text)
            if match:
                if match.group(1):
                    segments.append((match.group(1), False, False, False))
                formatted_text = match.group(3) or match.group(4)
                segments.append((formatted_text, True, True, False))
                remaining_text = remaining_text[match.end() :]
                continue

            # 太字: **text** or __text__
            match = re.match(r"^(.*?)(\*\*(.+?)\*\*|__(.+?)__)", remaining_text)
            if match:
                if match.group(1):
                    segments.append((match.group(1), False, False, False))
                formatted_text = match.group(3) or match.group(4)
                segments.append((formatted_text, True, False, False))
                remaining_text = remaining_text[match.end() :]
                continue

            # 斜体: *text* or _text_ (太字でない場合のみ)
            match = re.match(
                r"^(.*?)(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", remaining_text
            )
            if not match:
                match = re.match(r"^(.*?)(?<!_)_(?!_)(.+?)(?<!_)_(?!_)", remaining_text)
            if match:
                if match.group(1):
                    segments.append((match.group(1), False, False, False))
                segments.append((match.group(2), False, True, False))
                remaining_text = remaining_text[match.end() :]
                continue

            # 下線: <u>text</u>
            match = re.match(r"^(.*?)<u>(.+?)</u>", remaining_text)
            if match:
                if match.group(1):
                    segments.append((match.group(1), False, False, False))
                segments.append((match.group(2), False, False, True))
                remaining_text = remaining_text[match.end() :]
                continue

            # マッチしない場合は残り全てをプレーンテキストとして追加
            segments.append((remaining_text, False, False, False))
            break

        return segments if segments else [(text, False, False, False)]

    @staticmethod
    def apply_formatting_to_paragraph(paragraph, text: str, format_info: dict):
        """
        段落にテキストと書式を適用
        """
        if not text:
            return

        # セグメントを抽出
        segments = MarkdownParser.extract_formatting_segments(text)

        # 段落をクリア
        paragraph.text = ""

        # 各セグメントを追加
        for segment_text, is_bold, is_italic, is_underline in segments:
            if not segment_text:
                continue

            run = paragraph.add_run()
            run.text = segment_text

            if is_bold:
                run.font.bold = True
            if is_italic:
                run.font.italic = True
            if is_underline:
                run.font.underline = True


class PPTXGenerator:
    """PowerPointプレゼンテーション生成クラス"""

    def __init__(self):
        self.parser = MarkdownParser()

    def create_presentation(
        self, slides_data: Union[List[dict], List], output_path: str
    ) -> str:
        """
        マークダウン形式のスライドデータからPowerPointプレゼンテーションを生成

        Args:
            slides_data: スライド情報のリスト（Slideオブジェクトまたは辞書）
            output_path: 出力ファイルパス

        Returns:
            str: 生成されたファイルのパス
        """
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            for slide_item in slides_data:
                # Slideオブジェクトの場合
                if hasattr(slide_item, "content"):
                    content = getattr(slide_item, "content", "") or ""

                    # contentからタイトルとコンテンツを分離
                    slide_dict = self._parse_slide_content(content)
                else:
                    # 辞書の場合
                    slide_dict = {
                        "title": slide_item.get("title", "") or "",
                        "content": slide_item.get("content", "") or "",
                        "is_title_slide": slide_item.get("is_title_slide", False),
                    }

                self._add_slide(prs, slide_dict)

            # ファイルを保存
            prs.save(output_path)
            logger.info(f"Presentation generated successfully: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Presentation generation error: {str(e)}", exc_info=True)
            raise

    def _parse_slide_content(self, content: str) -> dict:
        """
        スライドのcontentからタイトルとコンテンツを分離
        各スライドでH1、H2は最初の1つのみを使用

        Args:
            content: Markdown形式のスライドコンテンツ

        Returns:
            dict: {'title': str, 'content': str, 'is_title_slide': bool}
        """
        lines = content.split("\n")
        title = ""
        remaining_lines = []
        is_title_slide = False
        h1_found = False
        h2_found = False

        for line in lines:
            line_stripped = line.strip()

            # H1を検出（タイトルスライド用）- 最初の1つのみ
            if line_stripped.startswith("# ") and not h1_found:
                title = line_stripped[2:].strip()
                h1_found = True
                is_title_slide = True
                continue

            # 2つ目以降のH1は警告してコンテンツに含める
            if line_stripped.startswith("# ") and h1_found:
                logger.warning(
                    f"Multiple H1 headings detected in slide. Subsequent H1 will be converted to H3: {line_stripped[:50]}"
                )
                # H3として扱う
                remaining_lines.append("### " + line_stripped[2:])
                continue

            # H2を検出（通常スライドのタイトル）- 最初の1つのみ、H1がない場合のみ
            if line_stripped.startswith("## ") and not h2_found and not h1_found:
                title = line_stripped[3:].strip()
                h2_found = True
                continue

            # 2つ目以降のH2は警告してコンテンツに含める
            if line_stripped.startswith("## ") and (h2_found or h1_found):
                logger.warning(
                    f"Multiple H2 headings detected in slide. Subsequent H2 will be converted to H3: {line_stripped[:50]}"
                )
                # H3として扱う
                remaining_lines.append("### " + line_stripped[3:])
                continue

            # その他の行はコンテンツに追加
            remaining_lines.append(line)

        return {
            "title": title,
            "content": "\n".join(remaining_lines),
            "is_title_slide": is_title_slide,
        }

    def _add_slide(self, prs: Presentation, slide_data: dict):
        """スライドを追加"""
        try:
            # レイアウトを選択
            if slide_data.get("is_title_slide", False):
                slide_layout = prs.slide_layouts[0]  # タイトルスライド
            else:
                slide_layout = prs.slide_layouts[1]  # タイトル+コンテンツ

            slide = prs.slides.add_slide(slide_layout)

            # タイトルの設定
            title_text = slide_data.get("title", "") or ""
            if slide.shapes.title and title_text:
                self._set_title(slide.shapes.title, title_text)

            # コンテンツの設定
            content = slide_data.get("content", "") or ""
            if content and len(slide.shapes) > 1:
                self._set_content(slide.shapes[1], content)

        except Exception as e:
            logger.error(f"Slide addition error: {str(e)}", exc_info=True)
            raise

    def _set_title(self, title_shape, title_text: str):
        """タイトルを設定"""
        if not title_text:
            title_text = ""

        title_shape.text = title_text
        if title_shape.text_frame.paragraphs:
            p = title_shape.text_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True

    def _set_content(self, content_shape, content: str):
        """コンテンツを設定"""
        if not content:
            return

        text_frame = content_shape.text_frame
        text_frame.clear()

        lines = content.split("\n")
        first_paragraph = True

        for line in lines:
            # 空行はスキップ
            if not line.strip():
                continue

            try:
                # マークダウンを解析
                parsed_text, indent_level, format_info = self.parser.parse_line(line)

                if not parsed_text:
                    continue

                # H1, H2はスキップ（タイトルとして既に処理済み）
                if format_info.get("is_heading", False):
                    heading_level = format_info.get("heading_level", 0)
                    if heading_level <= 2:
                        continue

                # 段落を追加
                if first_paragraph:
                    p = text_frame.paragraphs[0]
                    first_paragraph = False
                else:
                    p = text_frame.add_paragraph()

                # 見出しの場合（H3以降）
                if format_info.get("is_heading", False):
                    heading_level = format_info.get("heading_level", 1)

                    # 箇条書きを明示的に無効化
                    p.level = 0
                    if hasattr(p, "_element"):
                        # 箇条書きフォーマットを削除
                        pPr = p._element.get_or_add_pPr()
                        # buNoneを設定して箇条書きを無効化
                        if (
                            pPr.find(
                                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
                            )
                            is None
                        ):
                            from lxml import etree

                            etree.SubElement(
                                pPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone",
                            )

                    # テキストと書式を適用
                    self.parser.apply_formatting_to_paragraph(
                        p, parsed_text, format_info
                    )

                    # 見出し書式を適用
                    self._apply_heading_format(p, heading_level)

                # 番号付きリストの場合
                elif format_info.get("is_numbered", False):
                    p.level = min(indent_level, 8)  # PowerPointは最大9レベル（0-8）

                    # 番号付きリストの設定
                    if hasattr(p, "_element"):
                        from lxml import etree

                        pPr = p._element.get_or_add_pPr()

                        # 既存の箇条書き設定を削除
                        for elem in pPr.findall(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
                        ):
                            pPr.remove(elem)
                        for elem in pPr.findall(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
                        ):
                            pPr.remove(elem)

                        # 番号付きリストを設定
                        buAutoNum = pPr.find(
                            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum"
                        )
                        if buAutoNum is None:
                            buAutoNum = etree.SubElement(
                                pPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum",
                            )
                            buAutoNum.set("type", "arabicPeriod")  # 1. 2. 3. 形式
                            # 開始番号を設定
                            number = format_info.get("number", 1)
                            if number > 1:
                                buAutoNum.set("startAt", str(number))

                    # テキストと書式を適用
                    self.parser.apply_formatting_to_paragraph(
                        p, parsed_text, format_info
                    )

                    # 基本フォント設定
                    for run in p.runs:
                        if run.font.size is None:
                            run.font.size = Pt(14)

                # 箇条書きの場合
                elif format_info.get("is_bullet", False):
                    p.level = min(indent_level, 8)  # PowerPointは最大9レベル（0-8）

                    # テキストと書式を適用
                    self.parser.apply_formatting_to_paragraph(
                        p, parsed_text, format_info
                    )

                    # 基本フォント設定
                    for run in p.runs:
                        if run.font.size is None:
                            run.font.size = Pt(14)

                # 通常テキストの場合
                else:
                    # 箇条書きを明示的に無効化
                    p.level = 0
                    if hasattr(p, "_element"):
                        # 箇条書きフォーマットを削除
                        pPr = p._element.get_or_add_pPr()
                        # buNoneを設定して箇条書きを無効化
                        if (
                            pPr.find(
                                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
                            )
                            is None
                        ):
                            from lxml import etree

                            etree.SubElement(
                                pPr,
                                "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone",
                            )

                    # テキストと書式を適用
                    self.parser.apply_formatting_to_paragraph(
                        p, parsed_text, format_info
                    )

                    # 基本フォント設定
                    for run in p.runs:
                        if run.font.size is None:
                            run.font.size = Pt(14)

                p.space_after = Pt(6)

            except Exception as e:
                logger.error(f"Line processing error: {line[:50]}... - {str(e)}")
                # エラーが発生しても続行
                continue

    def _apply_heading_format(self, paragraph, heading_level: int):
        """見出し書式を適用（H3以降用）"""
        font_sizes = {3: 20, 4: 18, 5: 16, 6: 14}
        font_size = font_sizes.get(heading_level, 16)

        # 段落内の全てのrunに書式を適用
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = True


# 後方互換性のための関数
def create_presentation_from_markdown(slides_data: List[dict], output_path: str) -> str:
    """
    マークダウン形式のスライドデータからPowerPointプレゼンテーションを生成
    （後方互換性のための関数）
    """
    generator = PPTXGenerator()
    return generator.create_presentation(slides_data, output_path)
